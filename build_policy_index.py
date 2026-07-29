#!/usr/bin/env python3
"""
build_policy_index.py — rebuild policy_index.json for BTRauch/POA-HR-Compliance

Run from the root of a local clone of the repo:

    pip install pdfplumber python-docx openpyxl python-pptx
    python3 build_policy_index.py

Writes policy_index.json in the current directory (schema-compatible with v1.0,
so the hr-github-agent skill's existing search script keeps working unchanged).

Options:
    --dry-run          Report what would be indexed; write nothing.
    --limit N          Only process the first N documents (for a quick test).
    --out PATH         Output path (default: policy_index.json).
    --max-chars N      Skip files whose extracted text exceeds N chars (default 400000).
    --verbose          Print each file as it is processed.
"""

import argparse
import json
import os
import re
import sys
import time

# ---------------------------------------------------------------- config

INCLUDE_EXT = {'.pdf', '.docx', '.xlsx', '.pptx', '.doc', '.xls', '.txt', '.md'}

# Files/dirs never worth indexing
SKIP_DIR_PARTS = {'.git', '.github', '__pycache__'}
SKIP_NAME_PATTERNS = [
    r'_Error\.txt$',        # broken SharePoint export artifacts
    r'^~\$',                # Office lock files
    r'^\.',                 # dotfiles
]

# Chunking: aim near the existing index's ~2000-char median, split on headings
TARGET_CHUNK = 2000
MAX_CHUNK = 3800
MIN_CHUNK = 200          # merge anything smaller into the previous chunk

# Topic vocabulary — keep identical to the original so the skill's
# topic_filter values continue to work.
TOPIC_KEYWORDS = {
    'benefits':     ['benefit', 'insurance', 'fsa', 'flexible spending', 'dependent care',
                     '403b', '403(b)', 'retirement', 'pension', 'dental', 'vision',
                     'life insurance', 'disability', 'open enrollment', 'premium',
                     'health plan', 'medical plan', 'eap', 'workers compensation'],
    'leave':        ['leave', 'pto', 'paid time off', 'vacation', 'annual leave', 'sick',
                     'fmla', 'maternity', 'paternity', 'parental leave', 'bereavement',
                     'compassionate leave', 'sabbatical', 'time off'],
    'compensation': ['salary', 'compensation', 'payroll', 'wage', 'pay scale',
                     'salary matrix', 'salary scale', 'bonus', 'raise', 'merit',
                     'overtime', 'timesheet', 'grade'],
    'conduct':      ['code of conduct', 'conflict of interest', 'ethics', 'whistleblower',
                     'confidentiality', 'nda', 'non-disclosure', 'fraud', 'integrity',
                     'money laundering', 'trafficking'],
    'harassment':   ['harassment', 'sexual harassment', 'discrimination', 'psea',
                     'sexual exploitation', 'sea prevention', 'hostile', 'retaliation'],
    'travel':       ['travel', 'per diem', 'perdiem', 'lodging', 'airfare', 'mileage',
                     'itinerary', 'travel authorization', 'fly america', 'm&ie',
                     'trip report', 'advance request', 'incidental'],
    'expense':      ['expense', 'reimburse', 'receipt', 'invoice', 'payment request',
                     'credit card', 'affidavit', 'allowable', 'unallowable',
                     'cost policy', 'wire', 'ach', 'check request'],
    'onboarding':   ['onboarding', 'new hire', 'orientation', 'offer letter',
                     'recruitment', 'hiring', 'job description', 'interview',
                     'background check', 'i-9', 'probation'],
    'offboarding':  ['offboarding', 'termination', 'resignation', 'exit', 'separation',
                     'closeout', 'final paycheck', 'cobra'],
    'performance':  ['performance', 'evaluation', 'appraisal', 'review', 'smart goal',
                     '30-60-90', 'objectives', 'feedback', 'professional development'],
    'remote_work':  ['telecommut', 'remote work', 'work from home', 'flex work',
                     'flexible work', 'hybrid'],
    'compliance':   ['compliance', '2 cfr 200', 'cfr', 'audit', 'spot check', 'donor',
                     'federal', 'uei', 'sam.gov', 'debarment', 'suspension',
                     'section 508', 'section 889', 'nicra', 'risk assessment',
                     'subaward', 'subrecipient', 'visual compliance', 'accountability'],
    'child_safety': ['child safeguard', 'child protection', 'youth', 'minor',
                     'safeguarding'],
    'finance':      ['budget', 'accounting', 'gaap', 'chart of accounts', 'cash flow',
                     'financial report', 'cost share', 'in-kind', 'leverage', 'vat',
                     'tax', 'accrual', 'reclass', 'bank reconciliation', 'icr'],
    'procurement':  ['procurement', 'sole source', 'source selection', 'rfq', 'rfp',
                     'rfa', 'quotation', 'bid', 'vendor', 'purchase order',
                     'micro-purchase', 'competitive proposal', 'evaluation matrix',
                     'terms of reference', 'scope of work', 'contract', 'consultant'],
    'ai_policy':    ['artificial intelligence', 'generative ai', 'chatgpt', 'ai usage'],
    'forms':        ['form', 'template', 'checklist', 'request form', 'worksheet'],
    'training':     ['training', 'webinar', 'presentation', 'course', 'workshop',
                     'capacitaci'],
    'discipline':   ['disciplinary', 'discipline', 'corrective action', 'grievance',
                     'warning', 'misconduct', 'non-compliance', 'noncompliance'],
    'holiday':      ['holiday', 'holidays', 'observed'],
    'inventory':    ['inventory', 'equipment', 'asset', 'vehicle', 'disposition'],
    'it_security':  ['password', 'it acceptable use', 'privacy', 'personally identifying',
                     'pii', 'data protection', 'cyber'],
}

# Words that look like a heading line in extracted text
HEADING_RE = re.compile(
    r'^(?:'
    r'[A-Z][A-Z0-9 \-&/,\'()\.]{4,80}'           # ALL CAPS line
    r'|(?:\d+(?:\.\d+)*\.?\s+[A-Z][^\n]{3,80})'  # numbered section
    r'|(?:(?:Section|Annex|Appendix|Chapter|Article|Part)\s+[\dIVXA-Z][^\n]{0,70})'
    r')\s*$'
)


# ---------------------------------------------------------------- extractors

def extract_pdf(path):
    try:
        import pdfplumber
        out = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                t = page.extract_text() or ''
                if t.strip():
                    out.append(t)
        text = '\n'.join(out)
        if text.strip():
            return text
    except Exception:
        pass
    # fallback: poppler
    try:
        import subprocess
        r = subprocess.run(['pdftotext', '-layout', path, '-'],
                           capture_output=True, timeout=180)
        return r.stdout.decode('utf-8', 'ignore')
    except Exception:
        return ''


def extract_docx(path):
    try:
        import docx
    except ImportError:
        return ''
    try:
        doc = docx.Document(path)
        parts = []
        for p in doc.paragraphs:
            t = p.text.strip()
            if not t:
                continue
            style = (p.style.name or '') if p.style else ''
            # mark headings so the chunker can use them
            if style.startswith('Heading') or style == 'Title':
                parts.append('\n@@HEADING@@' + t)
            else:
                parts.append(t)
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                cells = [c for c in cells if c]
                if cells:
                    parts.append(' | '.join(dict.fromkeys(cells)))
        return '\n'.join(parts)
    except Exception:
        return ''


def extract_xlsx(path):
    try:
        import openpyxl
    except ImportError:
        return ''
    try:
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        parts = []
        for ws in wb.worksheets:
            parts.append('\n@@HEADING@@Sheet: ' + str(ws.title))
            rows = 0
            for row in ws.iter_rows(values_only=True):
                vals = [str(v).strip() for v in row if v is not None and str(v).strip()]
                if vals:
                    parts.append(' | '.join(vals))
                    rows += 1
                if rows > 400:      # spreadsheets can be enormous; cap per sheet
                    parts.append('[... additional rows omitted ...]')
                    break
        wb.close()
        return '\n'.join(parts)
    except Exception:
        return ''


def extract_pptx(path):
    try:
        from pptx import Presentation
    except ImportError:
        return ''
    try:
        prs = Presentation(path)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f'\n@@HEADING@@Slide {i}')
            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t:
                        parts.append(t)
                if getattr(shape, 'has_table', False):
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            parts.append(' | '.join(cells))
            if slide.has_notes_slide:
                n = slide.notes_slide.notes_text_frame.text.strip()
                if n:
                    parts.append('Notes: ' + n)
        return '\n'.join(parts)
    except Exception:
        return ''


def extract_plain(path):
    try:
        with open(path, encoding='utf-8', errors='ignore') as f:
            return f.read()
    except Exception:
        return ''


def extract_legacy(path):
    """Legacy .doc/.xls — try antiword/libreoffice if present, else skip."""
    import shutil, subprocess
    if path.lower().endswith('.doc') and shutil.which('antiword'):
        try:
            r = subprocess.run(['antiword', path], capture_output=True, timeout=120)
            return r.stdout.decode('utf-8', 'ignore')
        except Exception:
            pass
    return ''


EXTRACTORS = {
    '.pdf': extract_pdf, '.docx': extract_docx, '.xlsx': extract_xlsx,
    '.pptx': extract_pptx, '.txt': extract_plain, '.md': extract_plain,
    '.doc': extract_legacy, '.xls': extract_legacy,
}


# ---------------------------------------------------------------- processing

def clean(text):
    text = text.replace('\r\n', '\n').replace('\r', '\n')
    text = re.sub(r'[ \t]+', ' ', text)
    text = re.sub(r'\n{3,}', '\n\n', text)
    # drop lines that are only page furniture
    keep = []
    for line in text.split('\n'):
        s = line.strip()
        if re.fullmatch(r'(page\s*)?\d+(\s*(of|/)\s*\d+)?', s, re.I):
            continue
        keep.append(line)
    return '\n'.join(keep).strip()


def detect_topics(text, path):
    hay = (path + '\n' + text).lower()
    hits = []
    for topic, kws in TOPIC_KEYWORDS.items():
        score = sum(hay.count(k) for k in kws)
        if score:
            hits.append((score, topic))
    hits.sort(reverse=True)
    topics = [t for _, t in hits[:6]]
    return topics or ['general']


def chunk_text(text):
    """Split into ~TARGET_CHUNK pieces, preferring heading and paragraph breaks."""
    lines = text.split('\n')
    chunks, buf, heading = [], [], ''
    cur_heading = ''

    def flush():
        if not buf:
            return
        body = '\n'.join(buf).strip()
        if body:
            chunks.append({'heading': cur_heading, 'text': body})

    for line in lines:
        s = line.strip()
        explicit = s.startswith('@@HEADING@@')
        if explicit:
            s = s[len('@@HEADING@@'):].strip()
        is_heading = explicit or (len(s) <= 90 and bool(HEADING_RE.match(s)))

        if is_heading and sum(len(b) for b in buf) >= TARGET_CHUNK * 0.6:
            flush()
            buf = []
            cur_heading = s
            continue
        if is_heading and not buf:
            cur_heading = s
            continue

        buf.append(line)
        if sum(len(b) for b in buf) >= MAX_CHUNK:
            flush()
            buf = []

    flush()

    # merge undersized chunks forward
    merged = []
    for c in chunks:
        if merged and len(c['text']) < MIN_CHUNK:
            merged[-1]['text'] += '\n' + c['text']
        else:
            merged.append(c)
    for c in merged:
        c['char_count'] = len(c['text'])
    return merged


def should_skip(name):
    return any(re.search(p, name) for p in SKIP_NAME_PATTERNS)


def find_files(root):
    found = []
    for dirpath, dirnames, filenames in os.walk(root):
        dirnames[:] = [d for d in dirnames if d not in SKIP_DIR_PARTS]
        for fn in filenames:
            ext = os.path.splitext(fn)[1].lower()
            if ext not in INCLUDE_EXT or should_skip(fn):
                continue
            full = os.path.join(dirpath, fn)
            rel = os.path.relpath(full, root).replace(os.sep, '/')
            if rel.startswith('.'):
                continue
            found.append((rel, full))
    return sorted(found)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('--root', default='.')
    ap.add_argument('--out', default='policy_index.json')
    ap.add_argument('--dry-run', action='store_true')
    ap.add_argument('--limit', type=int)
    ap.add_argument('--max-chars', type=int, default=400000)
    ap.add_argument('--verbose', action='store_true')
    args = ap.parse_args()

    files = find_files(args.root)
    if args.limit:
        files = files[:args.limit]
    print(f'Found {len(files)} candidate documents under {os.path.abspath(args.root)}')
    if args.dry_run:
        from collections import Counter
        print(Counter(os.path.splitext(f)[1].lower() for f, _ in files))
        return 0

    documents, failed, skipped = [], [], []
    t0 = time.time()
    for i, (rel, full) in enumerate(files, 1):
        ext = os.path.splitext(rel)[1].lower()
        if args.verbose:
            print(f'  [{i}/{len(files)}] {rel}')
        elif i % 50 == 0:
            print(f'  ...{i}/{len(files)} ({time.time()-t0:.0f}s)')

        try:
            raw = EXTRACTORS[ext](full)
        except Exception as e:
            failed.append((rel, f'{type(e).__name__}: {e}'))
            continue

        text = clean(raw or '')
        if len(text) < 50:
            skipped.append((rel, 'no extractable text (scanned image or binary-only)'))
            continue
        if len(text) > args.max_chars:
            text = text[:args.max_chars]

        chunks = chunk_text(text)
        if not chunks:
            skipped.append((rel, 'no chunks produced'))
            continue

        documents.append({
            'path': rel,
            'filename': os.path.basename(rel),
            'folder': os.path.dirname(rel),
            'topics': detect_topics(text, rel),
            'total_chars': len(text),
            'chunk_count': len(chunks),
            'chunks': chunks,
        })

    index = {
        'version': '2.0',
        'repo': 'BTRauch/POA-HR-Compliance',
        'description': 'Pre-processed HR & policy document index for Partners of the Americas',
        'generated': time.strftime('%Y-%m-%d'),
        'documents': documents,
        'stats': {
            'total_documents': len(documents),
            'total_chunks': sum(d['chunk_count'] for d in documents),
            'total_characters': sum(d['total_chars'] for d in documents),
            'files_scanned': len(files),
            'no_text_extracted': len(skipped),
            'extraction_errors': len(failed),
        },
    }

    with open(args.out, 'w', encoding='utf-8') as f:
        json.dump(index, f, ensure_ascii=False, separators=(',', ':'))

    mb = os.path.getsize(args.out) / 1e6
    print(f'\nWrote {args.out}  ({mb:.1f} MB)')
    print(f'  documents indexed : {len(documents)}')
    print(f"  chunks            : {index['stats']['total_chunks']}")
    print(f"  characters        : {index['stats']['total_characters']:,}")
    print(f'  no text extracted : {len(skipped)}')
    print(f'  errors            : {len(failed)}')
    print(f'  elapsed           : {time.time()-t0:.0f}s')

    if skipped:
        with open('index_skipped.txt', 'w') as f:
            for r, why in skipped:
                f.write(f'{r}\t{why}\n')
        print('  -> see index_skipped.txt (candidates for OCR)')
    if failed:
        with open('index_errors.txt', 'w') as f:
            for r, why in failed:
                f.write(f'{r}\t{why}\n')
        print('  -> see index_errors.txt')
    return 0


if __name__ == '__main__':
    sys.exit(main())
Rebuild policy index across full repo (169 -> 1,441 documents)
